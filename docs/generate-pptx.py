"""
Generate: GitHub Copilot for Developers — Getting Started (detailed edition)
Outputs: docs/ghcp-getting-started.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ── Brand colours ──────────────────────────────────────────────────────────
GH_BLACK   = RGBColor(0x0D, 0x11, 0x17)   # GitHub dark bg
GH_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GH_BLUE    = RGBColor(0x23, 0x8B, 0xE6)   # Copilot accent blue
GH_GREY    = RGBColor(0x30, 0x36, 0x3D)   # Slide body bg
ACCENT_GRN = RGBColor(0x2D, 0xBA, 0x4E)   # GitHub green
ACCENT_YLW = RGBColor(0xF7, 0x8C, 0x00)   # Warning / highlight
ACCENT_PRP = RGBColor(0x79, 0x4E, 0xC8)   # Purple for skills
ACCENT_RED = RGBColor(0xCF, 0x22, 0x2E)   # Red for hooks/danger

W = Inches(13.333)   # widescreen 16:9
H = Inches(7.5)


# circled number markers for callouts
MARKERS = ["①","②","③","④","⑤","⑥","⑦","⑧","⑨","⑩"]

# ── Helpers ────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    return shape


def add_label(slide, text, l, t, w, h, size=20, bold=False, color=GH_WHITE,
              align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, color=GH_BLUE, height=Inches(0.06)):
    add_rect(slide, 0, 0, W, height, color)


def slide_number_tag(slide, num, total, color=GH_BLUE):
    add_label(slide, f"{num} / {total}",
              W - Inches(1.4), H - Inches(0.45), Inches(1.3), Inches(0.35),
              size=11, color=RGBColor(0x88,0x92,0x9A), align=PP_ALIGN.RIGHT)


# ── Slide builders ─────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle):
    layout = prs.slide_layouts[6]   # blank
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_BLACK)
    # left accent strip
    add_rect(sl, 0, 0, Inches(0.18), H, GH_BLUE)
    # copilot gradient panel
    add_rect(sl, Inches(0.18), 0, W - Inches(0.18), H, GH_BLACK)
    # decorative circle
    circ = sl.shapes.add_shape(9, W - Inches(4.2), Inches(-1.5), Inches(6), Inches(6))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x1A,0x22,0x2E)
    circ.line.fill.background()
    # title
    add_label(sl, title,
              Inches(0.6), Inches(1.8), Inches(9), Inches(1.8),
              size=44, bold=True, color=GH_WHITE)
    # blue underline
    add_rect(sl, Inches(0.6), Inches(3.55), Inches(2.2), Inches(0.07), GH_BLUE)
    # subtitle
    add_label(sl, subtitle,
              Inches(0.6), Inches(3.75), Inches(9.5), Inches(1.0),
              size=22, color=RGBColor(0xA0,0xAA,0xB4))
    # footer tag
    add_label(sl, "github.com/features/copilot",
              Inches(0.6), H - Inches(0.65), Inches(5), Inches(0.4),
              size=13, color=RGBColor(0x55,0x60,0x6A))
    return sl


def slide_section(prs, number, title, color=GH_BLUE):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_BLACK)
    add_rect(sl, 0, 0, Inches(0.18), H, color)
    add_label(sl, f"SECTION  {number:02d}", Inches(0.5), Inches(2.4), Inches(10), Inches(0.5),
              size=14, color=color)
    add_label(sl, title, Inches(0.5), Inches(2.9), Inches(11), Inches(1.4),
              size=40, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.5), Inches(4.35), Inches(1.6), Inches(0.06), color)
    return sl


def slide_content(prs, title, bullets, total, num,
                  bg=GH_GREY, note=None, code_block=None):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, bg)
    accent_bar(sl)
    add_label(sl, title,
              Inches(0.55), Inches(0.22), Inches(11.5), Inches(0.7),
              size=26, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.55), Inches(0.95), Inches(11.2), Inches(0.03), GH_BLUE)

    if code_block:
        # split area: bullets left, code right
        bw = Inches(5.4)
        add_bullets(sl, bullets, Inches(0.55), Inches(1.1), bw, Inches(5.6))
        add_rect(sl, Inches(6.2), Inches(1.05), Inches(6.8), Inches(5.7), GH_BLACK)
        add_label(sl, code_block,
                  Inches(6.35), Inches(1.2), Inches(6.55), Inches(5.4),
                  size=12, color=RGBColor(0xA8,0xCC,0x88), wrap=True)
    else:
        add_bullets(sl, bullets, Inches(0.55), Inches(1.1), Inches(12.2), Inches(5.6))

    if note:
        add_rect(sl, Inches(0.55), H - Inches(1.0), Inches(12.2), Inches(0.7),
                 RGBColor(0x1A,0x2A,0x3A))
        add_label(sl, f"💡  {note}",
                  Inches(0.7), H - Inches(0.98), Inches(12.0), Inches(0.62),
                  size=13, color=RGBColor(0xA0,0xC8,0xF0))
    slide_number_tag(sl, num, total)
    return sl


def add_bullets(slide, bullets, l, t, w, h):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        # item can be str  →  normal bullet
        #              (str, int) → (text, indent_level)
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)
        run = p.add_run()
        run.text = ("    " * level + ("• " if level == 0 else "◦ ")) + text
        run.font.size = Pt(17 if level == 0 else 15)
        run.font.bold = (level == 0)
        run.font.color.rgb = GH_WHITE if level == 0 else RGBColor(0xC0,0xCA,0xD4)


def slide_two_col(prs, title, left_head, left_items, right_head, right_items, total, num):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_GREY)
    accent_bar(sl)
    add_label(sl, title,
              Inches(0.55), Inches(0.22), Inches(11.5), Inches(0.7),
              size=26, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.55), Inches(0.95), Inches(11.2), Inches(0.03), GH_BLUE)

    col_w = Inches(5.8)
    # left column
    add_rect(sl, Inches(0.55), Inches(1.1), col_w, Inches(0.42), GH_BLUE)
    add_label(sl, left_head, Inches(0.65), Inches(1.12), col_w-Inches(0.2), Inches(0.38),
              size=15, bold=True)
    add_bullets(sl, left_items, Inches(0.55), Inches(1.6), col_w, Inches(5.0))

    # right column
    rx = Inches(6.95)
    add_rect(sl, rx, Inches(1.1), col_w, Inches(0.42), ACCENT_GRN)
    add_label(sl, right_head, rx+Inches(0.1), Inches(1.12), col_w-Inches(0.2), Inches(0.38),
              size=15, bold=True)
    add_bullets(sl, right_items, rx, Inches(1.6), col_w, Inches(5.0))

    slide_number_tag(sl, num, total)
    return sl


def slide_table(prs, title, headers, rows, total, num):
    from pptx.util import Inches, Pt
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_GREY)
    accent_bar(sl)
    add_label(sl, title,
              Inches(0.55), Inches(0.22), Inches(11.5), Inches(0.7),
              size=26, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.55), Inches(0.95), Inches(11.2), Inches(0.03), GH_BLUE)

    cols = len(headers)
    col_w = [Inches(v) for v in ([2.8] + [(13.333-0.55-0.55-2.8)/(cols-1)]*(cols-1))]
    row_h = Inches(0.52)
    top = Inches(1.1)

    table = sl.shapes.add_table(len(rows)+1, cols,
                                Inches(0.55), top,
                                Inches(12.23), row_h*(len(rows)+1)).table
    table.columns[0].width = col_w[0]
    for ci in range(1, cols):
        table.columns[ci].width = col_w[ci]

    def set_cell(cell, text, bg, fg=GH_WHITE, bold=False, size=14):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = fg

    for ci, h in enumerate(headers):
        set_cell(table.cell(0, ci), h, GH_BLUE, bold=True, size=14)
    for ri, row in enumerate(rows):
        bg = RGBColor(0x23,0x2B,0x33) if ri % 2 == 0 else RGBColor(0x1C,0x23,0x2A)
        for ci, val in enumerate(row):
            set_cell(table.cell(ri+1, ci), val, bg, size=13)

    slide_number_tag(sl, num, total)
    return sl


def slide_code(prs, title, description, code, total, num):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_GREY)
    accent_bar(sl)
    add_label(sl, title,
              Inches(0.55), Inches(0.22), Inches(11.5), Inches(0.7),
              size=26, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.55), Inches(0.95), Inches(11.2), Inches(0.03), GH_BLUE)
    add_label(sl, description,
              Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.7),
              size=16, color=RGBColor(0xA0,0xAA,0xB4))
    # code box
    add_rect(sl, Inches(0.55), Inches(1.75), Inches(12.2), Inches(5.3), GH_BLACK)
    add_label(sl, code,
              Inches(0.75), Inches(1.85), Inches(11.8), Inches(5.1),
              size=13, color=RGBColor(0xA8,0xCC,0x88), wrap=True)
    slide_number_tag(sl, num, total)
    return sl


def slide_closing(prs, total):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_BLACK)
    add_rect(sl, 0, 0, Inches(0.18), H, GH_BLUE)
    add_label(sl, "Start Today",
              Inches(0.6), Inches(1.6), Inches(10), Inches(1.2),
              size=46, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.6), Inches(2.85), Inches(2.0), Inches(0.07), GH_BLUE)
    steps = [
        "1.  Install GitHub Copilot extension in VS Code",
        "2.  Sign in with your GitHub account",
        "3.  Open Chat  →  try  /plan  with your next task",
        "4.  Use  /init  to generate workspace instructions for your repo",
        "5.  Run  /create-instructions  to capture team conventions",
        "6.  Explore  /create-prompt  or  /create-skill  for repeatable workflows",
    ]
    add_bullets(sl, steps, Inches(0.6), Inches(3.05), Inches(11.5), Inches(3.5))
    add_label(sl, "github.com/features/copilot  •  docs.github.com/copilot",
              Inches(0.6), H - Inches(0.6), Inches(11), Inches(0.4),
              size=13, color=RGBColor(0x55,0x60,0x6A))
    slide_number_tag(sl, total, total)
    return sl


def slide_annotated(prs, title, subtitle, code_text, callouts, total, num, accent=GH_BLUE):
    """
    Split slide: left = dark code panel with circled-number markers,
                 right = numbered callout bullets explaining each marker.
    code_text : plain string — include MARKERS[i] inline in the code lines.
    callouts  : list of str — one per marker, in order.
    """
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, GH_GREY)
    accent_bar(sl, color=accent)
    add_label(sl, title,
              Inches(0.55), Inches(0.22), Inches(11.5), Inches(0.7),
              size=24, bold=True, color=GH_WHITE)
    add_rect(sl, Inches(0.55), Inches(0.95), Inches(11.2), Inches(0.03), accent)
    if subtitle:
        add_label(sl, subtitle,
                  Inches(0.55), Inches(1.0), Inches(12.2), Inches(0.55),
                  size=14, color=RGBColor(0xA0,0xAA,0xB4))

    code_top = Inches(1.6) if subtitle else Inches(1.1)
    code_h   = H - code_top - Inches(0.3)

    # code panel (left ~55%)
    code_w = Inches(7.1)
    add_rect(sl, Inches(0.55), code_top, code_w, code_h, GH_BLACK)
    add_label(sl, code_text,
              Inches(0.72), code_top + Inches(0.12), code_w - Inches(0.3), code_h - Inches(0.2),
              size=12, color=RGBColor(0xA8,0xCC,0x88), wrap=True)

    # callout panel (right ~40%)
    cx = Inches(7.9)
    cw = Inches(5.1)
    for i, text in enumerate(callouts):
        ty = code_top + Inches(0.05) + i * Inches(0.72)
        # circled number badge
        badge = sl.shapes.add_shape(9, cx, ty + Inches(0.04), Inches(0.38), Inches(0.38))
        badge.fill.solid(); badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = MARKERS[i]
        br.font.size = Pt(11)
        br.font.bold = True
        br.font.color.rgb = GH_WHITE

        add_label(sl, text,
                  cx + Inches(0.46), ty, cw - Inches(0.46), Inches(0.68),
                  size=13, color=RGBColor(0xD0,0xD8,0xE4), wrap=True)

    slide_number_tag(sl, num, total)
    return sl


# ── Main ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    TOTAL = 39

    # ── 1  Title ────────────────────────────────────────────────────────────
    slide_title(prs,
        "GitHub Copilot\nfor Developers",
        "From zero to /plan — a practical deep-dive")

    # ── 2  Agenda ───────────────────────────────────────────────────────────
    slide_content(prs, "Agenda", [
        "What is GitHub Copilot?",
        "Built-in slash commands overview",
        "How /plan works — under the hood",
        "Customisation primitives — overview & decision guide",
        "Deep dive: Workspace Instructions",
        "Deep dive: File Instructions (.instructions.md)",
        "Deep dive: Prompts (.prompt.md)",
        "Deep dive: Skills (SKILL.md)",
        "Deep dive: Custom Agents (.agent.md)",
        "Deep dive: Hooks (.json)",
        "PoC walkthrough — OAuth2 + JWT auth system",
        "Best practices, security & getting-started checklist",
    ], TOTAL, 2)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 1 — What is GitHub Copilot?
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 1, "What is GitHub Copilot?")

    slide_content(prs, "GitHub Copilot — Overview", [
        "AI pair programmer embedded in VS Code, JetBrains, Visual Studio, and the CLI",
        "Powered by large language models (Claude Sonnet 4.6 in Copilot Chat)",
        "Two main surfaces:",
        ("Inline completions — next-token code suggestions as you type", 1),
        ("Copilot Chat — conversational agent with full tool use", 1),
        "Agent mode capabilities: read files, run commands, edit code, run tests — autonomously",
        "Everything the agent does is governed by your customisation files",
        "Session memory: agent tracks a plan across the whole conversation",
    ], TOTAL, 4,
    note="This session focuses on Copilot Chat in agent mode — the most powerful surface.")

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 2 — Slash Commands
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 2, "Built-in Slash Commands", color=ACCENT_GRN)

    slide_table(prs, "Built-in Slash Commands  ( type  /  in Chat )",
        ["Command", "Purpose", "What it Does"],
        [
            ["/plan",                "Research & planning",     "Interviews you with a questionnaire, then produces a phased implementation plan saved to session memory"],
            ["/init",                "Workspace bootstrap",     "Scans your codebase and auto-generates copilot-instructions.md with build commands, conventions, architecture notes"],
            ["/create-agent",        "Custom agent",            "Reviews conversation history, then guides creation of a .agent.md persona with restricted tools"],
            ["/create-instructions", "Coding conventions",      "Extracts patterns from your conversation and creates a persistent .instructions.md rule file"],
            ["/create-prompt",       "Reusable prompt",         "Packages a repeated task pattern into a .prompt.md slash command you can share with the team"],
            ["/create-skill",        "Workflow skill",          "Bundles a multi-step workflow + scripts/templates into a SKILL.md folder"],
            ["/create-hook",         "Lifecycle hook",          "Creates a .json hook that deterministically enforces policy at agent lifecycle events"],
        ],
        TOTAL, 6)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 3 — How /plan Works
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 3, "How  /plan  Works — Under the Hood", color=ACCENT_YLW)

    slide_annotated(prs,
        "The /plan Template File — Full Contents",
        "Location: assets/prompts/plan.prompt.md  inside the Copilot Chat VS Code extension",
        "---\nname: plan           " + MARKERS[0] + "\ndescription: Research and plan with the Plan agent  " + MARKERS[1] + "\nagent: Plan          " + MARKERS[2] + "\nargument-hint: Describe what you want to plan or research  " + MARKERS[3] + "\n---\n                     " + MARKERS[4] + "\nPlan my task.        " + MARKERS[5],
        [
            "name — the slug used as the /plan slash command",
            "description — shown in the command picker so users know what it does",
            "agent: Plan — routes this prompt to the built-in Plan agent mode (not the default agent)",
            "argument-hint — the placeholder text shown in the chat input box after you type /plan",
            "The --- closing the YAML frontmatter. Everything below is the prompt body.",
            "The entire prompt body is just this one sentence. All planning intelligence is inside the Plan agent, not here.",
        ],
        TOTAL, 8, accent=ACCENT_YLW)

    slide_content(prs, "What the Plan Agent Actually Does", [
        "The template is only a routing declaration — 7 lines total",
        "All intelligence lives inside the built-in Plan agent mode",
        "Step-by-step process when you type  /plan <your task>:",
        ("1. Parses the request — identifies domain, stack hints, implied requirements", 1),
        ("2. Identifies missing decisions → generates a targeted questionnaire", 1),
        ("3. Maps your answers to concrete technology choices and patterns", 1),
        ("4. Produces a dependency-ordered, sequenced implementation plan", 1),
        ("5. Saves the complete plan to session memory so it persists through implementation", 1),
        "Every downstream decision — naming, patterns, security — flows from your answers",
    ], TOTAL, 9,
    note="The agent knows what it doesn't know — it asks before it acts. Vague answers produce vague plans.")

    slide_two_col(prs,
        "PoC Walkthrough — /plan in Action",
        "The Questionnaire Answers",
        [
            "Framework?        .NET / ASP.NET Core",
            "OAuth provider?   Google",
            "OAuth flow?       Authorization Code Flow",
            "Token storage?    Redis",
            "Frontend?         React SPA",
            "Token strategy?   Access + Refresh tokens",
        ],
        "The 6-Phase Plan Produced",
        [
            "Phase 1 — Scaffold: .slnx, API project, test project, React frontend",
            "Phase 2 — Config: NuGet packages, JwtSettings, appsettings.json",
            "Phase 3 — Core services: TokenService (RS256) + RedisRefreshTokenStore (Lua)",
            "Phase 4 — Controller: 5 endpoints, secure cookie helpers",
            "Phase 5 — Frontend: AuthContext, axios interceptor, ProtectedRoute, pages",
            "Phase 6 — Security: headers middleware, rate limiter, CORS",
        ],
        TOTAL, 10)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 4 — Customisation Primitives Overview
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 4, "Customisation Primitives", color=GH_BLUE)

    slide_table(prs, "Which Primitive Should I Use?",
        ["Primitive", "File", "Location", "When to Use"],
        [
            ["Workspace Instructions", "copilot-instructions.md  or  AGENTS.md", ".github/  or root", "Always-on project-wide rules for every task"],
            ["File Instructions",      "*.instructions.md",         ".github/instructions/",  "Rules scoped to file types or specific task contexts"],
            ["Prompts",                "*.prompt.md",               ".github/prompts/",       "Single focused reusable task with parameterised inputs"],
            ["Hooks",                  "*.json",                    ".github/hooks/",         "Deterministic enforcement via shell at lifecycle events"],
            ["Custom Agents",          "*.agent.md",                ".github/agents/",        "Specialised persona with restricted tools / context isolation"],
            ["Skills",                 "SKILL.md  (in named folder)","github/skills/<name>/",  "Multi-step on-demand workflow with bundled scripts and assets"],
        ],
        TOTAL, 12)

    slide_content(prs, "Customisation Decision Shortcuts", [
        "Instructions vs Skill?",
        ("Applies to most / all work in the project  →  Workspace Instructions", 1),
        ("On-demand specialist task triggered by context  →  Skill", 1),
        "Skill vs Prompt?",
        ("Multi-step workflow with scripts, templates, or reference docs  →  Skill", 1),
        ("Single focused task with one clear output  →  Prompt", 1),
        "Skill vs Custom Agent?",
        ("Same tool set needed throughout all steps  →  Skill", 1),
        ("Need context isolation, or different tool restrictions per stage  →  Custom Agent", 1),
        "Hooks vs Instructions?",
        ("Instructions guide the agent  (non-deterministic)  →  Instructions", 1),
        ("Must block or enforce via shell commands regardless of model output  →  Hook", 1),
    ], TOTAL, 13)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 5 — Workspace Instructions
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 5, "Workspace Instructions", color=ACCENT_GRN)

    slide_annotated(prs,
        "Workspace Instructions — Template & Key Lines",
        "File: .github/copilot-instructions.md   (or AGENTS.md at repo root — pick one, not both)",
        "# Project Guidelines          " + MARKERS[0] + "\n\n## Code Style               " + MARKERS[1] + "\n- Use C# 12 primary constructors\n- 4-space indent, no tabs\n- Reference: src/SampleClass.cs\n\n## Architecture             " + MARKERS[2] + "\n- API layer: Controllers only route, no logic\n- Services own all business logic\n- Why: testability and separation\n\n## Build and Test           " + MARKERS[3] + "\n- Build:  dotnet build\n- Test:   dotnet test\n- Frontend: npm run dev\n\n## Conventions              " + MARKERS[4] + "\n- Never hardcode secrets — use dotnet user-secrets\n- All JWT must use RS256 — never HS256",
        [
            "No YAML frontmatter needed — this file is automatically applied to every chat request across the whole workspace.",
            "Code Style section: language/formatting rules. Link to exemplar files rather than duplicating code inline.",
            "Architecture section: component boundaries and the 'why' behind structural decisions. Helps agent avoid wrong patterns.",
            "Build and Test section: exact commands the agent will run automatically. Must be accurate — agent treats these as truth.",
            "Conventions section: project-specific rules that differ from defaults. If a linter enforces it, omit it here.",
        ],
        TOTAL, 15, accent=ACCENT_GRN)

    slide_two_col(prs,
        "Workspace Instructions — Principles & Anti-patterns",
        "Core Principles",
        [
            "Minimal by default — only what's relevant to every task",
            "Concise and actionable — every line should direct behavior",
            "Link, don't embed — reference docs instead of copying them",
            "Keep current — update when practices change",
            "Use AGENTS.md for monorepos: closer file in tree wins",
            "For large repos, use applyTo-based instructions for sub-areas (frontend, backend, tests)",
        ],
        "Anti-patterns to Avoid",
        [
            "Using BOTH copilot-instructions.md and AGENTS.md — pick one",
            "Kitchen-sink file: everything instead of what matters most",
            "Duplicating the README — link to it instead",
            "Obvious instructions like 'write clean code' — not actionable",
            "Stale instructions — agent will follow them even if outdated",
            "No build/test commands — agent then guesses and often gets it wrong",
        ],
        TOTAL, 16)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 6 — File Instructions
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 6, "File Instructions  (.instructions.md)", color=GH_BLUE)

    slide_annotated(prs,
        "File Instructions — Frontmatter & Key Lines",
        "File: .github/instructions/auth-patterns.instructions.md",
        "---\nname: auth-patterns          " + MARKERS[0] + "\ndescription: >              " + MARKERS[1] + "\n  Use when writing auth code,\n  JWT, OAuth2, or security\n  middleware. Covers token\n  validation and cookie rules.\napplyTo: \"**/Controllers/**\" " + MARKERS[2] + "\n---\n                             " + MARKERS[3] + "\n# Auth Patterns\n\n- Always use RS256 — never HS256  " + MARKERS[4] + "\n- Set HttpOnly, Secure, SameSite=Strict on all auth cookies\n- Validate issuer AND audience in JWT bearer middleware\n- Never log access tokens or refresh tokens",
        [
            "name — optional, defaults to filename. Shown in the instructions picker.",
            "description — the discovery surface. Write 'Use when...' with specific keywords. Agent reads this to decide if it's relevant. Keyword-rich is critical.",
            "applyTo — optional glob pattern. When a file matching this pattern is in context, this instruction auto-attaches. Omit for on-demand only.",
            "The --- closes YAML frontmatter. The body below is the instruction content.",
            "Every rule should be specific and actionable. The agent treats these as hard constraints when the instruction is loaded.",
        ],
        TOTAL, 18, accent=GH_BLUE)

    slide_table(prs, "File Instructions — Discovery Modes",
        ["Mode", "Trigger", "Frontmatter Required", "Best Use Case"],
        [
            ["On-demand",  "Agent detects task relevance from description keywords", "description  only",         "Task-based: migrations, refactoring, security, API design"],
            ["Explicit",   "Files matching applyTo glob are in the context window",  "applyTo  glob pattern",     "File-based: language standards, framework rules, test conventions"],
            ["Manual",     "User selects via Add Context > Instructions",            "Either or neither",         "Ad-hoc one-time attachment for a specific session"],
            ["Always-on",  "applyTo: '**'  matches every file",                     "applyTo: \"**\"  (caution!)", "Rare — only if truly universal; bloats every request context"],
        ],
        TOTAL, 19)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 7 — Prompts
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 7, "Prompts  (.prompt.md)", color=ACCENT_YLW)

    slide_annotated(prs,
        "Prompts — Frontmatter & Key Lines",
        "File: .github/prompts/generate-tests.prompt.md   →   invoked as  /generate-tests",
        "---\nname: generate-tests         " + MARKERS[0] + "\ndescription: >               " + MARKERS[1] + "\n  Generate xUnit test cases\n  for selected code. Covers\n  edge cases, error paths,\n  and existing test patterns.\nagent: agent                 " + MARKERS[2] + "\nargument-hint: >             " + MARKERS[3] + "\n  Paste the class or method\n  to test\ntools: [read, search, edit]  " + MARKERS[4] + "\nmodel: \"Claude Sonnet 4\"     " + MARKERS[5] + "\n---\n\nGenerate comprehensive xUnit tests for the selected code.  " + MARKERS[6] + "\n- Include happy path, edge cases, and error scenarios\n- Match patterns in existing test files: [tests](../../AuthSystem.Tests/)\n- Use descriptive [Fact] names in Given_When_Then format",
        [
            "name — the slug for the /generate-tests slash command. Must be filename-safe.",
            "description — shown in picker and used for agent auto-discovery. Include task keywords.",
            "agent: agent — runs in full agent mode with tool access. Options: agent, ask (chat only), plan, or a custom agent name.",
            "argument-hint — placeholder text shown in the chat input after you type the slash command.",
            "tools — restrict which tools this prompt can use. Principle of least privilege: only what's needed.",
            "model — override the default model for this specific prompt. Supports fallback array.",
            "Body: the actual task instructions. Reference files with Markdown links — agent will read them.",
        ],
        TOTAL, 21, accent=ACCENT_YLW)

    slide_two_col(prs,
        "Prompts — Principles & Anti-patterns",
        "Core Principles",
        [
            "Single task focus: one prompt = one well-defined output",
            "Output examples: show expected format when structure matters",
            "Reference files with Markdown links: [config](./config.json) — agent reads them",
            "Use argument-hint to tell users what input to provide",
            "Reuse over duplication: reference .instructions.md files instead of copying their content",
            "Both prompts and skills appear as / slash commands — skills for multi-step, prompts for single-task",
        ],
        "Anti-patterns to Avoid",
        [
            "Multi-task prompts: 'create and test and deploy' → split into separate prompts",
            "Vague descriptions: users won't know when to use it",
            "Over-tooling: giving a read-only task execute access",
            "Hardcoded paths that break on other machines",
            "No context references: agent guesses at file structure instead of reading it",
            "Using prompt when a Skill is better: if > 3 steps with assets → make a Skill",
        ],
        TOTAL, 22)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 8 — Skills
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 8, "Skills  (SKILL.md)", color=ACCENT_PRP)

    slide_annotated(prs,
        "Skills — SKILL.md Frontmatter & Key Lines",
        "File: .github/skills/run-e2e-tests/SKILL.md   (folder name must match 'name' field)",
        "---\nname: run-e2e-tests          " + MARKERS[0] + "\ndescription: >               " + MARKERS[1] + "\n  Run end-to-end tests using\n  Playwright. Use for verifying\n  frontend, debugging UI,\n  capturing failure screenshots.\nargument-hint: >             " + MARKERS[2] + "\n  Optionally specify a page\n  or feature to target\nuser-invocable: true         " + MARKERS[3] + "\ndisable-model-invocation: false  " + MARKERS[4] + "\n---\n\n# End-to-End Testing         " + MARKERS[5] + "\n\n## When to Use               " + MARKERS[6] + "\n- After every UI change\n- When debugging visual regressions\n\n## Procedure\n1. Start API:  dotnet run --project AuthSystem.API\n2. Run tests: [playwright script](./scripts/run-tests.js)  " + MARKERS[7] + "\n3. Screenshots saved to ./screenshots/",
        [
            "name — required, 1-64 chars, lowercase alphanumeric + hyphens. MUST match the folder name exactly or the skill silently fails.",
            "description — the 2-stage discovery surface: first ~100 tokens are read to decide relevance, then full body loads. Max 1024 chars.",
            "argument-hint — shown when user types /run-e2e-tests in chat, guiding what to type after the command.",
            "user-invocable: true — shows this skill as a /slash command. Set false to make it agent-only (auto-loaded, not user-triggered).",
            "disable-model-invocation: false — default. Set true to prevent automatic loading; only loads when user explicitly invokes.",
            "Body has no frontmatter boundary — everything after --- is the skill instructions loaded when triggered.",
            "'When to Use' section is critical — helps agent and user know exactly when this skill applies.",
            "Reference scripts with relative Markdown links. Agent loads them progressively only when needed — keeps context lean.",
        ],
        TOTAL, 24, accent=ACCENT_PRP)

    slide_content(prs, "Skills — Folder Structure & Progressive Loading", [
        "Recommended folder layout:",
        (".github/skills/<skill-name>/SKILL.md      ← required entry point", 1),
        (".github/skills/<skill-name>/scripts/      ← executable scripts", 1),
        (".github/skills/<skill-name>/references/   ← detailed docs loaded on demand", 1),
        (".github/skills/<skill-name>/assets/       ← templates, boilerplate files", 1),
        "Progressive loading — 3 stages:",
        ("Stage 1 — Discovery (~100 tokens): agent reads name + description to decide relevance", 1),
        ("Stage 2 — Instructions (<5000 tokens): loads full SKILL.md body when triggered", 1),
        ("Stage 3 — Resources: additional files load ONLY when explicitly referenced in body", 1),
        "Keep SKILL.md under 500 lines — move detail into references/ and link from the body",
        "Always use relative paths ( ./ ) for skill assets — absolute paths break portability",
    ], TOTAL, 25,
    note="Progressive loading keeps the context window lean. The agent only pays the cost of files it actually needs.")

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 9 — Custom Agents
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 9, "Custom Agents  (.agent.md)", color=ACCENT_YLW)

    slide_annotated(prs,
        "Custom Agents — Frontmatter & Key Lines",
        "File: .github/agents/code-reviewer.agent.md",
        "---\ndescription: >               " + MARKERS[0] + "\n  Use when: reviewing code for\n  security issues, OWASP Top 10,\n  code quality, or PR feedback.\nname: Code Reviewer          " + MARKERS[1] + "\ntools: [read, search]        " + MARKERS[2] + "\nmodel: \"Claude Sonnet 4\"     " + MARKERS[3] + "\nuser-invocable: true         " + MARKERS[4] + "\nagents: []                   " + MARKERS[5] + "\n---\n\nYou are a security-focused code reviewer.  " + MARKERS[6] + "\n\n## Constraints               " + MARKERS[7] + "\n- DO NOT edit or write any files\n- DO NOT run any shell commands\n- ONLY review and report findings\n\n## Approach\n1. Read all changed files\n2. Check each OWASP Top 10 category\n3. Return structured findings with severity",
        [
            "description — REQUIRED. The primary discovery surface for both the agent picker AND parent agents delegating to subagents. Use 'Use when:' pattern with specific trigger keywords.",
            "name — shown in the agent picker dropdown. Optional; defaults to filename without extension.",
            "tools — restrict to minimum needed. [read, search] = read-only; [] = conversational only; omit = all defaults.",
            "model — override the model. Supports fallback array: ['Claude Sonnet 4.5 (copilot)', 'GPT-5 (copilot)']",
            "user-invocable: true — shows in agent picker. Set false to hide from picker; agent still accessible as subagent.",
            "agents: [] — prevent this agent from invoking any subagents. Omit to allow all; list names to restrict.",
            "Persona statement in body — define the role clearly. This shapes the model's entire behavior.",
            "Constraints section — explicit DO NOT rules are critical for focused agents. Define exactly what it must NOT do.",
        ],
        TOTAL, 27, accent=ACCENT_YLW)

    slide_table(prs, "Custom Agent Tool Aliases",
        ["Alias", "What It Provides", "Use When"],
        [
            ["read",    "Read file contents",                    "Any agent that needs to inspect code or config"],
            ["edit",    "Edit / create files",                   "Implementation agents that write code"],
            ["search",  "Search files and text in workspace",    "Research or review agents"],
            ["execute", "Run shell commands in terminal",        "Build, test, or deploy agents — use carefully"],
            ["agent",   "Invoke other custom agents as subagents","Orchestrator agents that delegate to specialists"],
            ["web",     "Fetch URLs and web search",             "Agents that need external documentation or APIs"],
            ["todo",    "Manage the task list",                  "Long multi-step planning agents"],
            ["[]",      "No tools — conversational only",        "Q&A or explanation agents with no workspace access"],
        ],
        TOTAL, 28)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 10 — Hooks
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 10, "Hooks  (.json)", color=ACCENT_RED)

    slide_annotated(prs,
        "Hooks — Configuration Format & Key Lines",
        "File: .github/hooks/prevent-force-push.json   (team-shared)  or  .claude/settings.local.json  (local only)",
        "{\n  \"hooks\": {             " + MARKERS[0] + "\n    \"PreToolUse\": [      " + MARKERS[1] + "\n      {\n        \"type\": \"command\",  " + MARKERS[2] + "\n        \"command\": \"./scripts/validate-tool.sh\",  " + MARKERS[3] + "\n        \"windows\": \"pwsh ./scripts/validate-tool.ps1\",  " + MARKERS[4] + "\n        \"timeout\": 15       " + MARKERS[5] + "\n      }\n    ]\n  }\n}",
        [
            "hooks — top-level object. Contains one key per lifecycle event. Multiple events can be defined in the same file.",
            "PreToolUse — the event name. Hook fires BEFORE every tool call. Other events: PostToolUse, SessionStart, UserPromptSubmit, Stop, SubagentStart/Stop, PreCompact.",
            "type: 'command' — the only supported type. The hook always runs a shell command.",
            "command — the default command (runs on Linux/macOS). The script receives tool context as JSON on stdin.",
            "windows — platform override. Use this to provide a PowerShell equivalent on Windows. Also: linux and osx overrides available.",
            "timeout — seconds before the hook is killed. Keep hooks fast — long-running hooks block the agent.",
        ],
        TOTAL, 30, accent=ACCENT_RED)

    slide_annotated(prs,
        "Hooks — Input / Output Contract",
        "Hooks receive JSON on stdin describing the event. They return JSON on stdout to control agent behavior.",
        "# stdin received by hook script:\n{\n  \"tool_name\": \"run_in_terminal\",\n  \"tool_input\": { \"command\": \"git push --force\" }\n}\n\n# stdout returned by hook script:\n{                                   " + MARKERS[0] + "\n  \"hookSpecificOutput\": {           " + MARKERS[1] + "\n    \"hookEventName\": \"PreToolUse\",\n    \"permissionDecision\": \"deny\",  " + MARKERS[2] + "\n    \"permissionDecisionReason\":     " + MARKERS[3] + "\n      \"Force push is not allowed.\"\n  }\n}\n\n# Exit codes:                       " + MARKERS[4] + "\n# 0  = success (hook passed)\n# 2  = blocking error (stops the agent)\n# other = non-blocking warning",
        [
            "Return a JSON object on stdout. Fields: continue (bool), stopReason (str), systemMessage (str injected into context), hookSpecificOutput.",
            "hookSpecificOutput is used by PreToolUse and PostToolUse events to carry event-specific decisions back to the agent.",
            "permissionDecision controls whether the tool proceeds: 'allow' (proceed), 'ask' (prompt user), 'deny' (block the tool call).",
            "permissionDecisionReason is shown to the user when decision is 'ask' or 'deny'. Make it human-readable.",
            "Exit codes: 0 = hook passed normally, 2 = blocking error that stops agent execution, any other value = non-blocking warning logged but execution continues.",
        ],
        TOTAL, 31, accent=ACCENT_RED)

    slide_table(prs, "Hook Lifecycle Events",
        ["Event", "When It Fires", "Common Use Cases"],
        [
            ["SessionStart",      "First prompt of a new agent session",   "Inject project context, validate environment, greet user"],
            ["UserPromptSubmit",  "Every time user submits a prompt",       "Log prompts for audit, inject dynamic context"],
            ["PreToolUse",        "BEFORE every tool call",                 "Block dangerous commands (rm -rf, force push), require approval"],
            ["PostToolUse",       "AFTER successful tool completion",       "Auto-format files, run linter, enforce test pass after edit"],
            ["PreCompact",        "Before context window compaction",       "Save important state before summary replaces detail"],
            ["SubagentStart",     "When a subagent is invoked",             "Log subagent calls, inject subagent-specific context"],
            ["SubagentStop",      "When a subagent finishes",               "Validate subagent output, aggregate results"],
            ["Stop",              "Agent session ends",                     "Save session logs, cleanup temp files, send notifications"],
        ],
        TOTAL, 32)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 11 — PoC Walkthrough (already in /plan section, add detail)
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 11, "PoC — Real Obstacles & Resolutions", color=GH_BLUE)

    slide_table(prs, "Real Obstacles Encountered During Agent-Driven Implementation",
        ["Problem", "Root Cause", "Resolution"],
        [
            [".sln file not found after dotnet new sln",           ".NET 10 creates .slnx format, not .sln",                    "Use AuthSystem.slnx in all subsequent commands"],
            ["NuGet package version conflicts at build",           "Projects targeted net8.0 but SDK installed was .NET 10",     "Upgraded both .csproj TargetFramework to net10.0"],
            ["Authentication.Google package not found",            "No longer an implicit framework reference in .NET 10",       "Explicitly added as NuGet package reference"],
            ["npm --template flag not passed to create-vite",      "npm consumed the flag; never forwarded to the scaffolder",  "Used npx -y create-vite@latest authsystem-frontend -- --template react-ts"],
            ["300+ missing members compiling IDatabase fake",      "StackExchange.Redis v2.11 has a very large interface",       "Replaced hand-rolled class with Moq mock + targeted .Setup() calls"],
            ["Moq StringSetAsync setup never matched at runtime",  "Production code used implicit overload; Moq required exact match", "Changed production code to explicit 5-param overload; updated Moq setup to match"],
        ],
        TOTAL, 34)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 12 — Best Practices & Security
    # ══════════════════════════════════════════════════════════════════════
    slide_section(prs, 12, "Best Practices & Security", color=ACCENT_GRN)

    slide_two_col(prs,
        "Best Practices",
        "Effective Use of Copilot",
        [
            "Run /init on every new repo before writing code",
            "Answer /plan questionnaires completely — vague answers → vague plans",
            "Create .instructions.md for team conventions — every dev gets the same output",
            "Use principle of least privilege for tool lists in agents and prompts",
            "Review every generated file — treat it as a PR, not gospel",
            "Description quality = discoverability. Bad descriptions = ignored customisations.",
            "YAML frontmatter is unforgiving — quote values containing colons",
        ],
        "Common Mistakes",
        [
            "Committing appsettings.json with real secrets",
            "Using applyTo: '**' when content is only relevant to specific files",
            "Creating both copilot-instructions.md AND AGENTS.md — pick one",
            "Writing instructions so long the agent ignores the end",
            "Forgetting to push .github/ customisation files — team can't share them",
            "Hooks that run slowly and block every tool call",
            "Vague agent descriptions that prevent subagent delegation",
        ],
        TOTAL, 36)

    slide_content(prs, "Security Considerations with AI-Generated Code", [
        "Generated code is a starting point — always review for OWASP Top 10",
        "Copilot does not store or transmit your secrets — but bad prompts can embed them in code",
        "Always use dotnet user-secrets / environment variables — never hardcode credentials",
        "Key patterns to verify in any generated auth code:",
        ("JWT: algorithm pinned to RS256 (asymmetric) — never HS256 (symmetric shared secret)", 1),
        ("JWT: both issuer AND audience validated in bearer middleware configuration", 1),
        ("Cookies: HttpOnly=true, Secure=true, SameSite=Strict — all three required", 1),
        ("CORS: locked to known SPA origin — never wildcard *", 1),
        ("SQL / queries: always parameterised — no string concatenation", 1),
        ("Input validation: at system boundary only — not scattered through business logic", 1),
        "Write security rules into .instructions.md — agent will enforce them on every file",
    ], TOTAL, 37,
    note="This PoC generated RS256 JWT, HttpOnly cookies, Lua atomic token rotation, rate limiting, and security headers — still reviewed every file.")

    # ── Closing ─────────────────────────────────────────────────────────
    slide_closing(prs, TOTAL)

    out = r"c:\Projects\GHCP\docs\ghcp-getting-started.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({TOTAL} slides)")


if __name__ == "__main__":
    build()

